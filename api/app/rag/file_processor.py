import os
import logging
import uuid
from typing import List, Dict, Any
from fastapi import UploadFile, HTTPException
from app.config.config import settings

logger = logging.getLogger(__name__)


class FileProcessor:
    def __init__(self):
        self.supported_types = {
            "txt": self._process_text,
            "pdf": self._process_pdf,
            "docx": self._process_docx,
            "doc": self._process_doc,
            "pptx": self._process_pptx,
            "ppt": self._process_ppt,
            "xlsx": self._process_xlsx,
            "xls": self._process_xls,
        }

    async def save_file(self, file: UploadFile) -> Dict[str, Any]:
        """保存上传的文件"""
        try:
            # 获取文件扩展名
            file_ext = file.filename.split(".")[-1].lower() if "." in file.filename else ""

            if file_ext not in self.supported_types:
                raise HTTPException(status_code=400, detail=f"不支持的文件类型: {file_ext}")

            # 生成唯一文件名
            file_id = str(uuid.uuid4())
            filename = f"{file_id}.{file_ext}"
            file_path = os.path.join(settings.upload_dir, "rag_documents", filename)

            # 确保目录存在
            os.makedirs(os.path.dirname(file_path), exist_ok=True)

            # 保存文件
            contents = await file.read()
            with open(file_path, "wb") as f:
                f.write(contents)

            # 获取文件大小
            file_size = len(contents)

            return {
                "filename": file.filename,
                "file_path": file_path,
                "file_type": file_ext,
                "file_size": file_size
            }
        except Exception as e:
            logger.error(f"保存文件失败: {str(e)}")
            raise HTTPException(status_code=500, detail="文件保存失败")

    def process_file(self, file_path: str, file_type: str) -> List[Dict[str, Any]]:
        """处理文件并提取文本内容"""
        if file_type not in self.supported_types:
            raise ValueError(f"不支持的文件类型: {file_type}")

        try:
            processor = self.supported_types[file_type]
            text_content = processor(file_path)

            # 文本分块
            chunks = self._chunk_text(text_content)

            return chunks
        except Exception as e:
            logger.error(f"处理文件失败: {str(e)}")
            raise

    def _process_text(self, file_path: str) -> str:
        """处理文本文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def _process_pdf(self, file_path: str) -> str:
        """处理PDF文件"""
        try:
            import PyPDF2
            text = ""
            with open(file_path, "rb") as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except ImportError:
            logger.error("请安装PyPDF2: pip install PyPDF2")
            raise
        except Exception as e:
            logger.error(f"处理PDF失败: {str(e)}")
            raise

    def _process_docx(self, file_path: str) -> str:
        """处理DOCX文件"""
        try:
            import docx
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            logger.error("请安装python-docx: pip install python-docx")
            raise
        except Exception as e:
            logger.error(f"处理DOCX失败: {str(e)}")
            raise

    def _process_doc(self, file_path: str) -> str:
        """处理DOC文件（旧版Word文档）"""
        # 在实际应用中，可能需要使用antiword或其他工具
        # 这里简化处理，直接尝试转换为docx或返回错误
        logger.warning("DOC文件处理支持有限，建议转换为DOCX格式")
        try:
            # 尝试使用文本提取库
            import textract
            text = textract.process(file_path).decode("utf-8")
            return text
        except ImportError:
            logger.error("请安装textract: pip install textract")
            raise
        except Exception as e:
            logger.error(f"处理DOC失败: {str(e)}")
            raise

    def _process_pptx(self, file_path: str) -> str:
        """处理PPTX文件"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            logger.error("请安装python-pptx: pip install python-pptx")
            raise
        except Exception as e:
            logger.error(f"处理PPTX失败: {str(e)}")
            raise

    def _process_ppt(self, file_path: str) -> str:
        """处理PPT文件（旧版PowerPoint文档）"""
        logger.warning("PPT文件处理支持有限，建议转换为PPTX格式")
        try:
            import textract
            text = textract.process(file_path).decode("utf-8")
            return text
        except ImportError:
            logger.error("请安装textract: pip install textract")
            raise
        except Exception as e:
            logger.error(f"处理PPT失败: {str(e)}")
            raise

    def _process_xlsx(self, file_path: str) -> str:
        """处理XLSX文件"""
        try:
            import pandas as pd
            text = ""
            xl = pd.ExcelFile(file_path)
            for sheet_name in xl.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text += f"工作表: {sheet_name}\n"
                text += df.to_string() + "\n\n"
            return text
        except ImportError:
            logger.error("请安装pandas: pip install pandas")
            raise
        except Exception as e:
            logger.error(f"处理XLSX失败: {str(e)}")
            raise

    def _process_xls(self, file_path: str) -> str:
        """处理XLS文件（旧版Excel文档）"""
        try:
            import pandas as pd
            text = ""
            xl = pd.ExcelFile(file_path)
            for sheet_name in xl.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text += f"工作表: {sheet_name}\n"
                text += df.to_string() + "\n\n"
            return text
        except ImportError:
            logger.error("请安装pandas: pip install pandas")
            raise
        except Exception as e:
            logger.error(f"处理XLS失败: {str(e)}")
            raise

    def _chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        """将文本分块"""
        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = start + chunk_size
            if end > text_length:
                end = text_length

            # 确保不在单词中间分割
            if end < text_length:
                # 查找最近的分割点（句号、换行等）
                while end > start and text[end] not in ['.', '。', '!', '！', '?', '？', '\n']:
                    end -= 1

                # 如果没有找到合适的分割点，使用原始结束位置
                if end == start:
                    end = start + chunk_size

            chunk_text = text[start:end].strip()
            if chunk_text:  # 忽略空块
                chunks.append({
                    "text": chunk_text,
                    "start_index": start,
                    "end_index": end
                })

            # 移动到下一个块，考虑重叠
            start = end - overlap if end - overlap > start else end

        return chunks